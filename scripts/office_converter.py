import sys
from io import BytesIO

import fitz

def convert_to_word(pdf_path, output_path):

    doc = fitz.open(pdf_path)
    total_text = sum([len(page.get_text().strip()) for page in doc])

    if total_text < (50 * len(doc)):
        import pytesseract
        from PIL import Image
        from docx import Document

        doc_out = Document()
        for page in doc:
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            text = pytesseract.image_to_string(img)
            doc_out.add_paragraph(text)
            doc_out.add_page_break()

        doc_out.save(output_path)
    else:
        from pdf2docx import Converter
        cv = Converter(pdf_path)
        cv.convert(
            output_path,
            start=0, end=None,
            keep_page_layout=False, connected_border=True,
            line_overlap_margin=0.2, line_margin=0.2,
            word_margin=0.2, bottom_margin=5.0
        )
        cv.close()

def convert_to_excel(pdf_path, output_path):

    import pandas as pd
    import camelot
    import pdfplumber

    doc = fitz.open(pdf_path)
    total_pages = len(doc)
    doc.close()

    all_extracted_dfs = []

    for page_num in range(1, total_pages + 1):
        extracted_tables_dfs = []

        try:
            lattice_tables = camelot.read_pdf(pdf_path, pages=str(page_num), flavor="lattice")
            if len(lattice_tables) > 0 and any(t.df.shape[1] > 1 for t in lattice_tables):
                for t in lattice_tables:
                    if t.df.shape[1] > 1:
                        extracted_tables_dfs.append(t.df)
        except Exception:
            pass

        if not extracted_tables_dfs:
            try:
                stream_tables = camelot.read_pdf(pdf_path, pages=str(page_num), flavor="stream")
                if len(stream_tables) > 0 and any(t.df.shape[1] > 1 for t in stream_tables):
                    for t in stream_tables:
                        if t.df.shape[1] > 1:
                            extracted_tables_dfs.append(t.df)
            except Exception:
                pass

        # Attempt 3: pdfplumber structure matching
        if not extracted_tables_dfs:
            try:
                strategy_strict = {"vertical_strategy": "lines", "horizontal_strategy": "lines",
                                   "intersection_y_tolerance": 15}
                strategy_loose = {"vertical_strategy": "text", "horizontal_strategy": "text", "snap_tolerance": 5,
                                  "join_tolerance": 5}

                with pdfplumber.open(pdf_path) as pdf:
                    page = pdf.pages[page_num - 1]
                    table = page.extract_table(strategy_strict) or page.extract_table(strategy_loose)
                    if table:
                        df = pd.DataFrame([[cell for cell in row] for row in table if any(row)])
                        if df.shape[1] > 1:
                            extracted_tables_dfs.append(df)
            except Exception:
                pass

        if not extracted_tables_dfs:
            try:
                page_doc = fitz.open(pdf_path)
                page_text = page_doc[page_num - 1].get_text("text").strip()
                page_doc.close()

                if page_text:
                    lines = [line.strip() for line in page_text.split('\n') if line.strip()]
                    processed_rows = []
                    for line in lines:
                        if ',' in line:
                            processed_rows.append(line.split(','))
                        else:
                            processed_rows.append([line])

                    df = pd.DataFrame(processed_rows)
                    if not df.empty:
                        extracted_tables_dfs.append(df)
            except Exception:
                pass

        if extracted_tables_dfs:
            for df in extracted_tables_dfs:
                df.columns = range(df.shape[1])
                all_extracted_dfs.append(df)

    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
        if all_extracted_dfs:
            master_df = pd.concat(all_extracted_dfs, ignore_index=True)
            master_df.to_excel(writer, sheet_name='All Rows', index=False, header=False)
        else:
            pd.DataFrame(['No tabular data detected on any page.']).to_excel(writer, sheet_name='Result', index=False)


def convert_to_powerpoint(pdf_path, output_path):

    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    if len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    doc = fitz.open(pdf_path)

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        rect = page.rect

        prs.slide_width = Pt(rect.width)
        prs.slide_height = Pt(rect.height)

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        zoom = 2
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)

        img_bytes = pix.tobytes("png")
        image_stream = BytesIO(img_bytes)

        slide.shapes.add_picture(
            image_stream,
            0, 0,
            width=Pt(rect.width), height=Pt(rect.height)
        )

    prs.save(output_path)


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python office_converter.py <format> <input_pdf> <output_file>")
        sys.exit(1)

    target_format = sys.argv[1].lower()
    input_pdf = sys.argv[2]
    output_file = sys.argv[3]

    try:
        if target_format == "docx":
            convert_to_word(input_pdf, output_file)
        elif target_format == "xlsx":
            convert_to_excel(input_pdf, output_file)
        elif target_format == "pptx":
            convert_to_powerpoint(input_pdf, output_file)
        else:
            print(f"Unsupported format: {target_format}")
            sys.exit(1)

        print("SUCCESS")
    except Exception as e:
        import traceback
        print(f"ERROR: {str(e)}\n{traceback.format_exc()}", file=sys.stderr)
        sys.exit(1)